"""
Extracción de texto de documentos multimodales (RF-09).
Retorna lista de (numero_pagina, texto) para conservar metadatos de cita (RF-16).
"""
import io
from typing import Literal

DocType = Literal["pdf", "pptx", "txt"]


def extract_pages(file_bytes: bytes, doc_type: DocType) -> list[tuple[int, str]]:
    """Retorna lista de (pagina, texto). Para TXT la 'pagina' es siempre 1."""
    if doc_type == "pdf":
        return _extract_pdf(file_bytes)
    if doc_type == "pptx":
        return _extract_pptx(file_bytes)
    if doc_type == "txt":
        return _extract_txt(file_bytes)
    raise ValueError(f"Tipo de documento no soportado: {doc_type}")


def _extract_pdf(file_bytes: bytes) -> list[tuple[int, str]]:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append((i, text.strip()))
    return pages


def _extract_pptx(file_bytes: bytes) -> list[tuple[int, str]]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line.strip():
                        texts.append(line.strip())
        if texts:
            slides.append((i, "\n".join(texts)))
    return slides


def _extract_txt(file_bytes: bytes) -> list[tuple[int, str]]:
    text = file_bytes.decode("utf-8", errors="replace").strip()
    return [(1, text)] if text else []
